import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import re

def extract_all_pages(pdf_path):
    doc = fitz.open(pdf_path)
    pages_content = []
    for page_num, page in enumerate(doc):
        text = page.get_text()
        pages_content.append({
            'page_num': page_num,
            'text': text,
            'images': page.get_images()
        })
    doc.close()
    return pages_content

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            shape.text = subtitle
    return slide

def add_content_slide(prs, title, content_lines, layout_idx=1):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(5.0)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.space_after = Pt(8)
        p.level = 0
    
    return slide

def add_section_slide(prs, section_title):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = section_title
    return slide

def create_professional_presentation(pdf_path, output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    pages = extract_all_pages(pdf_path)
    
    add_title_slide(
        prs, 
        "基于2D CNN的AM/FM调制方式识别系统",
        "AM/FM Modulation Recognition System Based on 2D CNN"
    )
    
    add_content_slide(prs, "目录 | Contents", [
        "1. 项目概述与研究背景",
        "2. 系统架构设计",
        "3. 信号生成与调制原理",
        "4. 深度学习模型设计",
        "5. 实验结果与分析",
        "6. 总结与展望"
    ])
    
    add_section_slide(prs, "一、项目概述与研究背景")
    
    add_content_slide(prs, "研究背景 | Research Background", [
        "• 调制方式自动识别(AMR)是无线通信信号检测与分类的核心技术",
        "• AM（调幅）和FM（调频）是两种最基本的模拟调制方式",
        "• 传统方法依赖人工特征提取，需要领域专家知识，泛化能力有限",
        "• 深度学习方法可自动学习特征，无需人工设计",
        "• 2015年后，STFT+CNN做调制识别已成为非常成熟的通用方案"
    ])
    
    add_content_slide(prs, "主要贡献 | Main Contributions", [
        "[OK] 提出了基于频谱图的信号表示方法，将1D时间序列转换为2D图像",
        "[OK] 设计了高效的2D CNN模型架构，包含4层卷积和自适应池化",
        "[OK] 在AWGN信道、SNR>5dB的测试场景下，识别准确率稳定在95%以上",
        "[OK] 构建了交互式Web应用，支持实时信号生成和AI识别",
        "[OK] 为无线通信信号智能识别提供了新的思路"
    ])
    
    add_section_slide(prs, "二、系统架构设计")
    
    add_content_slide(prs, "整体架构 | Overall Architecture", [
        "信号流向:",
        "[原始信号] → [调制] → [加噪] → [频谱图转换] → [2D CNN] → [识别结果]",
        "",
        "• 信号生成: 产生AM/FM调制信号 → 1D时间序列",
        "• 加噪处理: 添加高斯白噪声 → 带噪1D信号",
        "• 频谱图转换: STFT时频分析 → 2D图像(时间×频率×功率)",
        "• CNN识别: 特征提取与分类 → AM/FM类别"
    ])
    
    add_content_slide(prs, "2D CNN模型架构 | Model Architecture", [
        "输入: 频谱图 [1, 17, 61] (通道, 频率, 时间)",
        "",
        "Conv2d(32) + ReLU + MaxPool2d",
        "Conv2d(64) + ReLU + MaxPool2d",
        "Conv2d(128) + ReLU + MaxPool2d",
        "Conv2d(256) + ReLU + AdaptiveAvgPool2d",
        "Flatten → Linear(512) → Dropout → Linear(2)",
        "",
        "输出: [AM概率, FM概率]"
    ])
    
    add_section_slide(prs, "三、信号生成与调制原理")
    
    add_content_slide(prs, "AM调制原理 | AM Modulation", [
        "AM调制公式: s(t) = [1 + m·cos(2πfmt)] · cos(2πfct)",
        "",
        "• m: 调制指数 (通常 0 < m ≤ 1)",
        "• fm: 消息信号频率",
        "• fc: 载波频率",
        "• 载波幅度随消息信号变化",
        "• 广泛应用于广播、电视等领域"
    ])
    
    add_content_slide(prs, "FM调制原理 | FM Modulation", [
        "FM调制公式: s(t) = sin(2πfct + β·cos(2πfmt))",
        "",
        "• β: 调制指数 (调制深度)",
        "• fm: 消息信号频率",
        "• fc: 载波频率",
        "• 载波瞬时频率随消息信号变化",
        "• 具有更好的抗噪声性能"
    ])
    
    add_section_slide(prs, "四、深度学习模型设计")
    
    add_content_slide(prs, "模型配置 | Model Configuration", [
        "• 优化器: Adam (学习率 0.001)",
        "• 损失函数: CrossEntropyLoss",
        "• 训练轮次: 20 epochs",
        "• 批处理大小: 64",
        "• 数据集划分: 80% 训练, 20% 测试",
        "• 正则化: Dropout (p=0.5)",
        "• 数据增强: 添加不同强度的高斯噪声"
    ])
    
    add_section_slide(prs, "五、实验结果与分析")
    
    add_content_slide(prs, "实验设置 | Experimental Setup", [
        "• 采样频率: 1000 Hz",
        "• 信号时长: 1秒",
        "• 载波频率: 100 Hz",
        "• 消息频率: 5 Hz",
        "• SNR范围: -10dB 到 20dB",
        "• 每种调制类型生成1000个样本",
        "• 总共2000个训练样本"
    ])
    
    add_content_slide(prs, "性能评估 | Performance Evaluation", [
        "• SNR > 5dB时，识别准确率稳定在95%以上",
        "• 低SNR条件下，准确率有所下降但仍保持较高水平",
        "• 模型对AM和FM两种调制方式均具有良好的识别能力",
        "• 混淆矩阵显示模型在两类之间的误判率较低",
        "• 训练曲线显示模型快速收敛，损失函数稳定下降"
    ])
    
    add_section_slide(prs, "六、总结与展望")
    
    add_content_slide(prs, "总结 | Summary", [
        "• 成功实现了基于2D CNN的AM/FM调制方式识别系统",
        "• 将1D时间序列信号转换为2D频谱图进行特征提取",
        "• 模型在多种SNR条件下均表现出色",
        "• 构建了完整的信号处理和深度学习流程",
        "• 验证了深度学习方法在调制识别任务上的有效性"
    ])
    
    add_content_slide(prs, "未来工作 | Future Work", [
        "• 扩展到更多调制类型（QAM, PSK等）",
        "• 优化模型架构以提高低SNR性能",
        "• 探索更先进的深度学习模型（Transformer, Attention）",
        "• 实现实时信号处理和识别",
        "• 部署到实际无线通信系统中进行验证"
    ])
    
    add_title_slide(prs, "谢谢！", "Thank You!")
    
    prs.save(output_path)
    print(f"演示文稿已成功创建: {output_path}")
    print(f"共创建 {len(prs.slides)} 张幻灯片")

if __name__ == "__main__":
    pdf_file = "AM_FM_Modulation_Recognition_Paper.pdf"
    output_file = "AM_FM_Presentation.pptx"
    create_professional_presentation(pdf_file, output_file)
